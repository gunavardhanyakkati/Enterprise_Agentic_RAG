import logging
from pathlib import Path
from typing import Optional

from src.exceptions import PDFParsingException, PDFValidationError
from src.schemas.pdf_parser.models import PaperSection, ParserType, PdfContent

logger = logging.getLogger(__name__)


class OfficeParser:
    """Office document parser for .docx, .pptx, .xlsx files."""

    def __init__(self):
        """Initialize Office parser."""
        self._parsers_available = self._check_parsers()
        logger.info(f"Office parser initialized with availability: {self._parsers_available}")

    def _check_parsers(self) -> dict:
        """Check which Office parser libraries are available."""
        availability = {
            "docx": False,
            "pptx": False,
            "xlsx": False,
        }

        try:
            import docx

            availability["docx"] = True
        except ImportError:
            logger.warning("python-docx not installed - .docx parsing disabled")

        try:
            import pptx

            availability["pptx"] = True
        except ImportError:
            logger.warning("python-pptx not installed - .pptx parsing disabled")

        try:
            import openpyxl

            availability["xlsx"] = True
        except ImportError:
            logger.warning("openpyxl not installed - .xlsx parsing disabled")

        return availability

    async def parse_document(self, file_path: Path) -> Optional[PdfContent]:
        """Parse an Office document and extract content.

        Args:
            file_path: Path to the Office document (.docx, .pptx, .xlsx)

        Returns:
            PdfContent object with extracted text and structure, or None if parsing failed
        """
        if not file_path.exists():
            logger.error(f"Office document not found: {file_path}")
            raise PDFValidationError(f"Office document not found: {file_path}")

        file_extension = file_path.suffix.lower()
        logger.info(f"Parsing Office document: {file_path.name} ({file_extension})")

        try:
            if file_extension == ".docx" and self._parsers_available["docx"]:
                return await self._parse_docx(file_path)
            elif file_extension == ".pptx" and self._parsers_available["pptx"]:
                return await self._parse_pptx(file_path)
            elif file_extension in [".xlsx", ".xls"] and self._parsers_available["xlsx"]:
                return await self._parse_xlsx(file_path)
            else:
                logger.error(f"Unsupported file type or parser not available: {file_extension}")
                return None

        except Exception as e:
            logger.error(f"Failed to parse Office document {file_path}: {e}")
            raise PDFParsingException(f"Office document parsing failed: {e}")

    async def _parse_docx(self, file_path: Path) -> PdfContent:
        """Parse a Word document (.docx)."""
        import docx
        from docx.document import Document as DocxDocument

        doc: DocxDocument = docx.Document(file_path)

        sections = []
        current_section = {"title": "Content", "content": "", "level": 1}

        # Extract text from paragraphs, grouping by heading style
        for para in doc.paragraphs:
            if para.style.name.startswith("Heading"):
                # Save previous section if it has content
                if current_section["content"].strip():
                    sections.append(PaperSection(**current_section))

                # Determine heading level
                level = 1
                if "Heading" in para.style.name:
                    try:
                        level = int(para.style.name.split()[-1])
                    except (ValueError, IndexError):
                        level = 1

                current_section = {
                    "title": para.text.strip(),
                    "content": "",
                    "level": level,
                }
            else:
                # Add paragraph text to current section
                if para.text.strip():
                    current_section["content"] += para.text + "\n\n"

        # Add final section
        if current_section["content"].strip():
            sections.append(PaperSection(**current_section))

        # Get full text
        full_text = "\n\n".join([para.text for para in doc.paragraphs if para.text.strip()])

        return PdfContent(
            sections=sections,
            figures=[],
            tables=[],
            raw_text=full_text,
            references=[],
            parser_used=ParserType.DOCLING,  # TODO: Create Office-specific parser type
            metadata={
                "source": "office_parser",
                "file_type": "docx",
                "paragraph_count": len(doc.paragraphs),
                "section_count": len(sections),
            },
        )

    async def _parse_pptx(self, file_path: Path) -> PdfContent:
        """Parse a PowerPoint presentation (.pptx)."""
        from pptx import Presentation

        prs = Presentation(file_path)

        sections = []
        full_text = []

        # Each slide becomes a section
        for idx, slide in enumerate(prs.slides, 1):
            slide_title = f"Slide {idx}"
            slide_content = ""

            # Extract title from first title shape if present
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    if shape.shape_type == 14:  # Title placeholder
                        slide_title = shape.text
                    else:
                        slide_content += shape.text + "\n\n"

            if slide_content.strip():
                sections.append(
                    PaperSection(
                        title=slide_title,
                        content=slide_content.strip(),
                        level=1,
                    )
                )
                full_text.append(f"# {slide_title}\n\n{slide_content}")

        return PdfContent(
            sections=sections,
            figures=[],
            tables=[],
            raw_text="\n\n".join(full_text),
            references=[],
            parser_used=ParserType.DOCLING,  # TODO: Create Office-specific parser type
            metadata={
                "source": "office_parser",
                "file_type": "pptx",
                "slide_count": len(prs.slides),
                "section_count": len(sections),
            },
        )

    async def _parse_xlsx(self, file_path: Path) -> PdfContent:
        """Parse an Excel spreadsheet (.xlsx)."""
        import openpyxl
        from openpyxl.worksheet.worksheet import Worksheet

        wb = openpyxl.load_workbook(file_path, data_only=True)

        sections = []
        full_text = []

        # Each worksheet becomes a section
        for sheet_name in wb.sheetnames:
            ws: Worksheet = wb[sheet_name]
            sheet_content = ""

            # Get cell values in a readable format
            for row in ws.iter_rows(values_only=True):
                # Filter out None and empty rows
                row_values = [str(val) for val in row if val is not None]
                if row_values:
                    sheet_content += "\t".join(row_values[:20]) + "\n"  # Limit to 20 columns

            if sheet_content.strip():
                sections.append(
                    PaperSection(
                        title=sheet_name,
                        content=sheet_content.strip(),
                        level=1,
                    )
                )
                full_text.append(f"# {sheet_name}\n\n{sheet_content}")

        return PdfContent(
            sections=sections,
            figures=[],
            tables=[],
            raw_text="\n\n".join(full_text),
            references=[],
            parser_used=ParserType.DOCLING,  # TODO: Create Office-specific parser type
            metadata={
                "source": "office_parser",
                "file_type": "xlsx",
                "sheet_count": len(wb.sheetnames),
                "section_count": len(sections),
            },
        )
